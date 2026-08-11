from pathlib import Path

from PIL import Image
from pptx import Presentation

from mp4slides.config import OutputConfig
from mp4slides.models import DetectedSlide
from mp4slides.pdf_writer import _fit_or_paginate_text, _register_font, write_pdf
from mp4slides.pptx_writer import write_pptx


def make_image(path: Path) -> None:
    image = Image.new("RGB", (640, 360), (240, 240, 240))
    image.save(path)


def make_slide(image_path: Path) -> DetectedSlide:
    text = "\u65e5\u672c\u8a9e\u306e\u6587\u5b57\u8d77\u3053\u3057"
    return DetectedSlide(
        start=0.0,
        end=5.0,
        representative_time=4.5,
        image_bytes=b"",
        analysis_gray=None,
        analysis_edges=None,
        image_path=str(image_path),
        transcript=text,
    )


def test_pptx_contains_notes(tmp_path: Path) -> None:
    image_path = tmp_path / "slide.png"
    make_image(image_path)
    output = tmp_path / "out.pptx"
    slide = make_slide(image_path)
    write_pptx([slide], output, OutputConfig(format="pptx"))
    presentation = Presentation(output)
    notes = presentation.slides[0].notes_slide.notes_text_frame.text
    assert slide.transcript in notes


def test_pptx_notes_can_replace_newlines_with_spaces(tmp_path: Path) -> None:
    image_path = tmp_path / "slide.png"
    make_image(image_path)
    output = tmp_path / "out.pptx"
    slide = make_slide(image_path)
    slide.transcript = "first line\nsecond line"
    config = OutputConfig(format="pptx", pptx_notes_newline_mode="space")
    write_pptx([slide], output, config)
    presentation = Presentation(output)
    notes = presentation.slides[0].notes_slide.notes_text_frame.text
    assert "first line second line" in notes
    assert "first line\nsecond line" not in notes


def test_pdf_writes_japanese_transcript(tmp_path: Path) -> None:
    image_path = tmp_path / "slide.png"
    make_image(image_path)
    output = tmp_path / "out.pdf"
    write_pdf([make_slide(image_path)], output, OutputConfig(format="pdf"))
    assert output.exists()
    assert output.stat().st_size > 1000


def test_side_by_side_pdf_keeps_all_text_across_pages() -> None:
    config = OutputConfig(format="pdf", pdf_transcript_mode="side-by-side")
    font_name = _register_font(config)
    text = "あ" * 5000
    _, _, pages = _fit_or_paginate_text(
        text,
        font_name,
        config.pdf_font_size,
        config.pdf_min_font_size,
        width=120.0,
        height=120.0,
    )
    assert len(pages) > 1
    assert "".join(line for page in pages for line in page) == text
