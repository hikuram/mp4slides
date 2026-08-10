from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from .config import OutputConfig
from .models import DetectedSlide


def _fit_size(image_width: int, image_height: int, box_width: int, box_height: int) -> tuple[int, int]:
    scale = min(box_width / image_width, box_height / image_height)
    return int(round(image_width * scale)), int(round(image_height * scale))


def _image_dimensions(path: str | Path) -> tuple[int, int]:
    from PIL import Image

    with Image.open(path) as image:
        return image.size


def write_pptx(slides: list[DetectedSlide], output_path: str | Path, config: OutputConfig) -> Path:
    target = Path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(config.slide_width_in)
    presentation.slide_height = Inches(config.slide_height_in)
    blank_layout = presentation.slide_layouts[6]

    for index, item in enumerate(slides, start=1):
        if not item.image_path:
            raise ValueError(f"Slide {index} has no image_path")
        slide = presentation.slides.add_slide(blank_layout)
        width_px, height_px = _image_dimensions(item.image_path)
        fit_width, fit_height = _fit_size(
            width_px,
            height_px,
            presentation.slide_width,
            presentation.slide_height,
        )
        left = int((presentation.slide_width - fit_width) / 2)
        top = int((presentation.slide_height - fit_height) / 2)
        slide.shapes.add_picture(item.image_path, left, top, width=fit_width, height=fit_height)
        notes = slide.notes_slide.notes_text_frame
        if notes is not None:
            notes.text = (
                f"[{item.start:.3f} - {item.end:.3f}]\n"
                f"{item.transcript}".rstrip()
            )

    presentation.save(target)
    return target
